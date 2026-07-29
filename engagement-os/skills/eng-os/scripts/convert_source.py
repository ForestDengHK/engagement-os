#!/usr/bin/env python3
"""Convert one source document to faithful, citable Markdown for the reference pack.

Deterministic extraction only: text + tables with per-page/slide anchors, plus image
extraction to a sibling images dir.

The converter only drops what bytes can prove decorative; everything else survives for triage.
A raster on *most* pages/slides is template furniture (logo, footer mark) and one under a few KB
is an icon or rule — both are counted and reported, so the pass stays auditable. A raster repeated
on just a few units is NOT dropped: decks legitimately re-show a content diagram (build slides, a
recap), so the first copy is kept and only the duplicate files go. A diagram drawn in vector
shapes is no raster at all — invisible to image extraction — so a PDF page full of vector art is
rendered as a page snapshot, and a pptx slide full of shape diagrams is rendered via LibreOffice.
Everything that survives is tagged `[uncertain]` for the agent to classify
[decorative]/[content] and OCR — that part is a judgment + vision task. See `eng-ingest-source`.

EXTRACTION IS DELEGATED, NOT REIMPLEMENTED. docx goes through `pandoc` and pptx through
`markitdown` — the same tools the `docx` / `pptx` skills read with. They settled document order,
heading levels, lists, footnotes, nested tables and shape coverage years ago; hand-rolling that
on top of python-docx/python-pptx just reproduces their bugs. pdf stays on pymupdf and xlsx on
openpyxl because those need per-page / per-sheet control to place the anchors, which a
whole-file converter cannot give.

What this script owns is the PACKAGING the pack's discipline depends on and no general converter
provides: the provenance header (source path, md5, unit counts), the citable anchor per unit,
image extraction to disk with decorative auto-drop, and one uniform interface across six formats.
Each output names its extractor, and a fallback says plainly what it lost.

Still reach for the matching SKILL, per document, when the script's own output falls short:
tracked changes, comments, chart data labels, or a scanned PDF needing real OCR.

Usage:
    python convert_source.py <source_path> [--out <md_path>] [--images-dir <dir>]

Supported: .pdf .pptx .ppsx .docx .xlsx .csv .png .jpg .jpeg .html .htm
Dependencies are imported lazily; a missing one degrades that format with a clear message
rather than crashing. Install as needed: pip install pymupdf python-pptx python-docx openpyxl
(PEP-668/Homebrew Python: add --user --break-system-packages, or use a venv.)

A sourced-from-the-web document has an origin URL that the file itself does not carry, and a
research claim has to cite something retrievable. Pass `--source-url` and it lands in the
provenance header next to the path and the md5.
"""
import argparse
import hashlib
import os
import pathlib
import re
import sys
import datetime as _dt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import figure_contract as fc


def md5(path):
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


# Set from --source-url. A document we downloaded has an origin the file does not record, and
# a research claim must cite something a reader can retrieve — so it belongs in the header,
# beside the path and the md5, for every format and not just HTML.
SOURCE_URL = None


def header(src, extra=""):
    origin = f"> **Origin URL:** {SOURCE_URL}  \n" if SOURCE_URL else ""
    return (
        f"# {os.path.basename(src)}\n\n"
        f"> **Source:** `{src}`  \n"
        f"{origin}"
        f"> **Converted:** {_dt.date.today().isoformat()}  \n"
        f"> **md5:** `{md5(src)}`  \n"
        f"{extra}\n---\n\n"
    )


MIN_IMG_BYTES = 6 * 1024        # below this is an icon, bullet, rule or spacer
MIN_IMG_PIXELS = 120            # narrowest side; logos and dividers fall under this
FURNITURE_MIN_COPIES = 4        # fewer than this can never be ruled template furniture
FURNITURE_UNIT_SHARE = 0.5      # …and it must sit on at least half the units
MIN_VECTOR_DRAWINGS = 20        # a page drawing this many vector shapes is a diagram, not text
COVER_RASTER_SHARE = 0.5        # a raster covering this much of the page IS the page's visual
LEGIBLE_LONG_EDGE = 700         # under this a diagram's own labels are too small to read
MIN_DIAGRAM_SHAPES = 10         # a slide composing this many drawn shapes is a diagram


def img_prefix(src):
    """Namespace for one document's extracted images inside a shared pack images dir."""
    stem = os.path.splitext(os.path.basename(src))[0]
    return re.sub(r"[^A-Za-z0-9._-]+", "_", stem) + "__"


class ImageCollector:
    """Collects extracted images, dropping the ones that are decorative by construction.

    Two heuristics do almost all the work, and both are decidable from the bytes:

    * **Frequency.** A raster on most pages/slides (at least half, and at least four copies)
      is template furniture — a logo, a footer mark. Below that bar an image might be a
      content diagram a deck legitimately shows twice, so the FIRST copy survives for triage
      and only the duplicate files are removed (identical bytes — no information lost).
      Deleting a unique image is a judgment call, and the collector only makes calls bytes can
      make. Found on the real AIB deck: a 93KB architecture diagram shown on two consecutive
      slides was auto-deleted as "template furniture" before triage ever saw it.
    * **Size.** Under a few KB, or narrower than ~120px, it is an icon or a rule.

    Everything dropped is counted and reported, so the pass stays auditable: the lossless rule
    is about not losing *information*, and a logo carries none. What survives is what a human
    would actually have to look at.

    Links are written relative to **the markdown file**, not to the images dir's parent. Those
    two coincide only when the images dir is exactly `<md_dir>/images`; every other layout the
    skill actually instructs (`_md/images/<topic>/`) produced links that resolved nowhere.

    Names are namespaced by source document. The per-unit name a converter generates
    (`p6_img1.png`) is unique only *within* one document, but a reference pack points every
    document at ONE shared images dir — so the eleventh document's page 6 quietly overwrote the
    first document's page 6. The markdown still rendered, still had a plausible caption stub, and
    now showed a different report's figure under our citation. Found on the Deloitte research
    E2E: 27 names claimed by 2-5 documents each. Prefixing with the source stem makes the
    filename mean what the citation says it means.
    """

    def __init__(self, images_dir, link_base=None, prefix=""):
        self.dir = images_dir
        self.link_base = link_base or os.path.dirname(images_dir.rstrip("/"))
        self.prefix = prefix
        self.kept = []                  # (relpath, unit) in emit order
        self._by_digest = {}            # digest -> [(path, rel, unit), ...]
        self._names = {}                # emitted filename -> relpath (for inline placement)
        self.dropped_small = 0
        self.consolidated = 0            # raster components retained inside a page snapshot
        self.stale_removed = 0            # old files from an earlier conversion of this source

    def add(self, blob, unit, name):
        """Offer one image. Returns True if written to disk."""
        if len(blob) < MIN_IMG_BYTES:
            self.dropped_small += 1
            return False
        name = self.prefix + os.path.basename(name)
        digest = hashlib.md5(blob).hexdigest()
        path = os.path.join(self.dir, name)
        os.makedirs(self.dir, exist_ok=True)
        with open(path, "wb") as f:
            f.write(blob)
        rel = os.path.relpath(path, self.link_base)
        self._by_digest.setdefault(digest, []).append((path, rel, unit))
        return True

    def finalise(self, total_units=None):
        """Decide each unique image's fate. Returns (kept, dropped_furniture, dropped_dup).

        Three outcomes per digest:
        * furniture (>= FURNITURE_MIN_COPIES and on >= half the units) → delete every copy;
        * a low-frequency repeat → keep the first copy, delete the duplicate files (same
          bytes, so nothing is lost — the survivor is placed at its first occurrence);
        * unique → keep.
        `total_units` is the page/slide count where the caller knows it (pdf, pptx); None
        (docx) falls back to the absolute-copy bar alone.
        """
        dropped_furniture = 0
        dropped_dup = 0
        for copies in self._by_digest.values():
            furniture = (len(copies) >= FURNITURE_MIN_COPIES and
                         (total_units is None or
                          len(copies) >= total_units * FURNITURE_UNIT_SHARE))
            if furniture:
                for path, _rel, _unit in copies:
                    try:
                        os.remove(path)
                    except OSError:
                        pass
                dropped_furniture += len(copies)
                continue
            _path, rel, unit = copies[0]
            self.kept.append((rel, unit))
            self._names[os.path.basename(rel)] = rel
            for dup_path, _rel, _unit in copies[1:]:
                try:
                    os.remove(dup_path)
                except OSError:
                    pass
                dropped_dup += 1
        self.kept.sort(key=lambda t: t[1])
        # Regeneration can legitimately change the representation: for example a converter
        # upgrade replaces twelve diagram fragments with one readable page snapshot. Leaving
        # the old names in the shared images directory makes them look like live evidence even
        # though no markdown references them. Clean only this source's namespaced files, and
        # only after the new conversion has completed its decisions.
        live = {os.path.abspath(os.path.join(self.link_base, rel)) for rel, _unit in self.kept}
        if self.prefix and os.path.isdir(self.dir):
            for name in os.listdir(self.dir):
                path = os.path.join(self.dir, name)
                if name.startswith(self.prefix) and os.path.isfile(path) \
                        and os.path.abspath(path) not in live:
                    try:
                        os.remove(path)
                        self.stale_removed += 1
                    except OSError:
                        pass
        return self.kept, dropped_furniture, dropped_dup

    def link_for(self, name):
        """Relative link for a kept image by filename, or None if it was dropped.
        `finalise()` must have run — decorative-drop is only decidable across the whole doc.

        Callers hold the un-prefixed name they generated, while `_names` is keyed by the
        namespaced one, so try both rather than making every call site know about the prefix.
        """
        base = os.path.basename(name)
        return self._names.get(self.prefix + base) or self._names.get(base)


IMG_TOKEN = fc.IMG_TOKEN
_IMG_TOKEN_RE = fc.IMG_TOKEN_RE


def img_block(rel, alt=""):
    """A placed image plus the caption stub that keeps the placement honest.

    Shape and markers come from `figure_contract` — `triage_images.py` rewrites exactly this
    block and `eng_lint.py` counts exactly these markers, so all three must read one definition.
    """
    block = fc.caption_stub_block(rel)
    label = alt.strip()
    if label:                                    # markitdown's accessibility description
        block = block.replace(f"![{os.path.basename(rel)}]({rel})", f"![{label}]({rel})", 1)
    return block


def place_inline_images(md, coll):
    """Resolve the inline placement tokens the converters emitted.

    An image only means something where it sat in the document — a scoring table or an
    architecture diagram listed at the bottom of the file under "triage these" has lost the
    clause it belonged to. But whether an image survives the decorative filter is only decidable
    once the whole document has been read, so the converters emit a token in place and this
    resolves it afterwards: kept → a markdown image at its original position, dropped → gone.
    A token may carry alt text after a `|` (markitdown's accessibility description, often the
    only written account of a figure) — it becomes the figure's label.
    """
    def sub(m):
        name, _, alt = m.group(1).partition("|")
        rel = coll.link_for(name)
        return img_block(rel, alt) if rel else ""
    return _IMG_TOKEN_RE.sub(sub, md)


def img_section(collector, total_units=None):
    if isinstance(collector, list):                  # legacy call-site safety
        kept, dropped_furniture, dropped_dup, dropped_small, consolidated = (
            [(i, 0) for i in collector], 0, 0, 0, 0)
    else:
        kept, dropped_furniture, dropped_dup = collector.finalise(total_units)
        dropped_small = collector.dropped_small
        consolidated = collector.consolidated
    note = ""
    if dropped_furniture or dropped_dup or dropped_small:
        bits = []
        if dropped_furniture:
            bits.append(f"{dropped_furniture} template furniture (logo / footer mark on most units)")
        if dropped_dup:
            bits.append(f"{dropped_dup} duplicate copies (identical bytes — first occurrence kept)")
        if dropped_small:
            bits.append(f"{dropped_small} under {MIN_IMG_BYTES // 1024}KB (icon / rule / spacer)")
        note = f"\n{fc.AUTODROP_LINE} {' · '.join(bits)}. " \
               "Furniture and icons carry no information, and duplicates survive at their first " \
               "occurrence — everything else is below.\n"
    if not kept:
        return ("\n---\n\n## Images\n" + note + "\nNo content-bearing images.\n") if note else ""
    # Careful with the wording: the lint counts occurrences of the caption tag, so this preamble
    # must describe the stub without spelling it — boilerplate that names its own marker makes
    # every converted file permanently fail the rule it exists to serve.
    # The count is the other half of the audit: after triage the surviving figures are whatever
    # the ledger says they are, and without a record of how many arrived, a document whose
    # figures were all deleted reads exactly like a document that never had any. That is how
    # the AIB deck lost 13 diagrams in silence. `eng_lint.rule_images_accounted` reconciles it.
    consolidation_note = (
        f"\n{fc.CONSOLIDATED_LINE} {consolidated} embedded raster component(s); each complete "
        "page is retained once at readable resolution, so its circles, arrows and fragments "
        "are not separate triage items.\n" if consolidated else "")
    lines = [f"\n---\n\n{fc.HEADING_PENDING}\n\n{fc.EXTRACTED_LINE} {len(kept)}\n",
             note, consolidation_note,
             "Each is also **placed inline** at the position it held in the source, under a "
             "caption stub. This index is the checklist: classify each `[decorative]` (delete "
             "the file, the inline block and this line) / `[content]` (keep + write the caption) "
             "/ `[uncertain]` (OCR inline, then retag `[ocr-done]`). Do not delete an "
             "`[uncertain]` image before its text is captured. A diagram, chart or figure is "
             "content even when it recurs across pages or sits under a brand background — only "
             "logos, icons, rules and bare backgrounds are decorative.\n"]
    for im, _unit in kept:
        lines.append(f"- `[{fc.UNTRIAGED}]` ![{os.path.basename(im)}]({im})")
    return "\n".join(lines) + "\n"


def _tool(name):
    """Path to an external converter, or None."""
    from shutil import which
    return which(name)


def _run(cmd):
    import subprocess
    try:
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        return r.stdout if r.returncode == 0 and r.stdout.strip() else None
    except Exception:
        return None


MAX_HEADING_CHARS = 120     # above this it is a mis-styled paragraph, not a heading


def _plain_heading(title):
    """Strip the emphasis wrapper pandoc faithfully carries into a heading.

    A real Word heading arrives as `**4.5. DECLARATION...**` or `<u>5.2. COST EVALUATION</u>`
    because the style bolded it. Keeping the wrapper leaves `## Section 42: **SCOPE**` in the
    output and makes the clause detection read the `**` instead of the number.
    """
    t = title.strip()
    t = re.sub(r"^(?:(?:\*\*|__|<u>|\[)\s*)+", "", t, flags=re.I)
    t = re.sub(r"(?:\s*(?:\*\*|__|</u>|\]))+$", "", t, flags=re.I)
    return t.strip()


def number_headings(md, label="Section"):
    """Give every markdown heading a stable, citable anchor, keeping its level.

    The extraction tools return the document's real heading hierarchy; what they can't give is
    an anchor a downstream claim can cite. A buyer's native clause number is the strongest
    anchor because it survives pagination and matches what an evaluator sees: `## 2.1 Timetable`
    becomes `## §2.1 Timetable`. Only headings without a native number get the synthetic
    fallback (`## Section 12: Minimum Requirements`). The outline hierarchy is never flattened.

    Three things the first version got wrong, all visible on the real GNI RFT:
    * the synthetic counter counted the `§`-numbered headings too, so the fallback sequence
      arrived full of holes (`Section 50` → `§4.5` → `Section 53`) and read like data loss;
    * an empty styled paragraph became a bare `#` — a heading with no text and no anchor;
    * a whole paragraph styled as Heading 3 in Word became `### Section 21: • The Contracting
      Entity may withhold payments pursuant to...`, which is body text wearing an anchor. A
      heading longer than `MAX_HEADING_CHARS` is demoted back to a paragraph; the text is kept,
      only the false anchor goes.
    """
    out, anchors, synthetic = [], 0, 0
    native = re.compile(
        r"^(?:(?:section|clause)\s+)?"
        r"(?P<num>\d{1,3}(?:\.\d{1,3}){0,5})"
        r"(?:\s*[.):\-–—]\s*|\s+)(?P<title>.+)$",
        re.I)
    for line in md.splitlines():
        m = re.match(r"^(#{1,6})\s+(.*)$", line)
        if not m:
            out.append(line)
            continue
        title = _plain_heading(m.group(2))
        if not title:
            continue                                 # empty styled paragraph — drop, don't anchor
        if len(title) > MAX_HEADING_CHARS:
            out.append(title)                        # mis-styled paragraph — keep text, drop anchor
            continue
        anchors += 1
        if title.startswith("§"):
            out.append(f"{m.group(1)} {title}")
            continue
        clause = native.match(title)
        if clause:
            out.append(f"{m.group(1)} §{clause.group('num')} {clause.group('title').strip()}")
        else:
            synthetic += 1
            out.append(f"{m.group(1)} {label} {synthetic}: {title}")
    return "\n".join(out), anchors


_TOC_NESTED = re.compile(r"\[([^\[\]]*?)\s*\[(\d+)\]\(#[^)]*\)\]\(#[^)]*\)")
_INTERNAL_LINK = re.compile(r"\[([^\[\]]*)\]\(#[^)]*\)")


def flatten_internal_links(md):
    """Turn Word's exported table of contents back into readable text.

    Pandoc renders a TOC field as a link whose own label contains another link —
    `[1. IMPORTANT INFORMATION [4](#important-information)](#important-information)`. That is
    not valid nesting in GFM, so it renders as literal brackets, and a hundred lines of it is
    the first thing anyone reading the converted RFT sees. The anchors are worthless anyway:
    `number_headings` rewrites the heading ids they point at. Keep the text and the page number,
    drop the link machinery. External links are untouched.
    """
    md = _TOC_NESTED.sub(lambda m: f"{m.group(1)} — p.{m.group(2)}", md)
    return _INTERNAL_LINK.sub(lambda m: m.group(1), md)


def _pptx_shapes(container):
    """Yield shapes depth-first: a grouped diagram keeps its labels in child shapes, and
    `slide.shapes` only yields the top level — those labels would be lost silently."""
    for sh in container:
        if getattr(sh, "shape_type", None) == 6:      # GROUP
            yield from _pptx_shapes(sh.shapes)
        else:
            yield sh


def _pptx_diagram_slides(prs):
    """Slides whose figure is COMPOSED of drawn shapes (org chart, SmartArt, flow diagram)
    rather than placed as a picture. python-pptx reads each shape's text but cannot draw the
    composition, so without a rendered snapshot the figure — usually the slide's whole point —
    leaves no image behind. Found on the real GNI_Team_Structure deck: 61 auto shapes per
    slide, zero pictures extracted, the reporting lines gone.
    """
    needed = []
    for i, slide in enumerate(prs.slides, 1):
        n = 0
        for sh in _pptx_shapes(slide.shapes):
            st = getattr(sh, "shape_type", None)
            if st in (1, 3, 5):                       # AUTO_SHAPE, CHART, FREEFORM — drawn, not placed
                n += 1
            elif st == 19:                            # GRAPHIC_FRAME: a table is extracted as
                try:                                  # text; a chart or SmartArt diagram is not
                    if sh.has_chart or "diagram" in sh._element.graphic.graphicData.uri:
                        n += 1
                except Exception:
                    pass
        if n >= MIN_DIAGRAM_SHAPES:
            needed.append(i)
    return needed


def _pptx_slide_snapshots(src, slides_needed, coll):
    """Render shape-diagram slides to images via LibreOffice, the same renderer the pptx
    skill uses. Returns {slide_no: name} for snapshots written. Only the needed slides are
    rendered; a slide whose figure is a placed picture already has its image extracted.
    """
    out = {}
    if not slides_needed or not _tool("soffice"):
        return out
    import shutil
    import subprocess
    import tempfile
    tmp = tempfile.mkdtemp(prefix="engos-slides-")
    try:
        r = subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                            "--outdir", tmp, src],
                           capture_output=True, text=True, timeout=300)
        pdf = os.path.join(tmp, os.path.splitext(os.path.basename(src))[0] + ".pdf")
        if r.returncode != 0 or not os.path.exists(pdf):
            return out
        import fitz
        doc = fitz.open(pdf)
        for i in sorted(slides_needed):
            if i > doc.page_count:
                continue
            pix = doc[i - 1].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            name = f"s{i}_slide.png"
            if coll.add(pix.tobytes("png"), i, name):
                out[i] = name
        doc.close()
    except Exception:
        return out
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    return out


def _snapshot_warning(missed):
    return (f"\n> **{missed} slide(s) hold shape diagrams that were NOT rendered** — no "
            "LibreOffice (`soffice`) on this machine, and python-pptx cannot draw composed "
            "shapes. Install LibreOffice and re-convert, or read those slides from the "
            "original file.\n")


def convert_pdf(src, images_dir, link_base=None):
    try:
        import fitz  # pymupdf
    except ImportError:
        return None, "pymupdf not installed — run: pip install pymupdf (PEP-668: add --user --break-system-packages)"
    doc = fitz.open(src)
    out = [header(src, f"> **Pages:** {doc.page_count}  ")]
    coll = ImageCollector(images_dir, link_base, img_prefix(src))
    for i, page in enumerate(doc, 1):
        out.append(f"## Page {i}:\n")
        text = page.get_text("text").strip()
        out.append(text + "\n" if text else "_(no extractable text — likely a diagram/image page; see extracted images)_\n")
        covered = False
        widest = 0
        rasters = []
        for j, img in enumerate(page.get_images(full=True)):
            try:
                pix = fitz.Pixmap(doc, img[0])
                if pix.n - pix.alpha >= 4:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                if min(pix.width, pix.height) < MIN_IMG_PIXELS:
                    coll.dropped_small += 1
                    continue
                try:
                    # a raster already covering most of the page IS the page's visual —
                    # a snapshot on top would only duplicate it
                    area = sum(r.get_area() for r in page.get_image_rects(img[0]))
                    if area >= page.rect.get_area() * COVER_RASTER_SHARE:
                        covered = True
                except Exception:
                    pass
                name = f"p{i}_img{j}.png"
                rasters.append((pix.tobytes("png"), name))
                widest = max(widest, pix.width, pix.height)
            except Exception:
                continue
        # A diagram drawn in vector shapes is no embedded raster at all — get_images() cannot
        # see it, and the page's text is only its scattered labels. Render the page itself so
        # the figure survives. Found on the real AIB deck: the roadmap and architecture pages
        # (68–229 vector drawings each) produced zero images and no hint a figure was there.
        # A second case the vector rule misses: the page's figures WERE extracted, but every
        # one is a thumbnail whose own labels are unreadable at that size. Extraction
        # "succeeded" and the reader still cannot see what the diagram says — on the AIB deck
        # a triage pass had to caption one such figure "read the full-size figure in the
        # source PDF", which is exactly the trip back to the binary the md-first rule exists
        # to avoid. The page itself carries them at readable size, so render it.
        thumbnails_only = widest and widest < LEGIBLE_LONG_EDGE
        snapshot_added = False
        if not covered:
            try:
                if len(page.get_drawings()) >= MIN_VECTOR_DRAWINGS or thumbnails_only:
                    snap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    name = f"p{i}_page.png"
                    if coll.add(snap.tobytes("png"), i, name):
                        out.append(IMG_TOKEN % name)
                        snapshot_added = True
            except Exception:
                pass
        if snapshot_added:
            # The snapshot is the complete content figure. PDF slide decks routinely encode
            # one diagram as a dozen raster/vector fragments; emitting those fragments beside
            # the page turns one judgment into a dozen and loses the relationships between
            # them. The full page contains every component at readable resolution and in
            # context, so retain it once and record how many components it subsumed.
            coll.consolidated += len(rasters)
        else:
            for data, name in rasters:
                if coll.add(data, i, name):
                    out.append(IMG_TOKEN % name)     # resolved after decorative filtering
    tail = img_section(coll, doc.page_count)         # runs finalise(); must precede placement
    return place_inline_images("\n".join(out), coll) + tail, None


def convert_pptx(src, images_dir, link_base=None):
    try:
        from pptx import Presentation
    except ImportError:
        return None, "python-pptx not installed — run: pip install python-pptx (PEP-668: add --user --break-system-packages)"
    prs = Presentation(src)
    coll = ImageCollector(images_dir, link_base, img_prefix(src))

    # markitdown owns the text: it reaches shape types python-pptx makes you hunt for and
    # carries each picture's embedded accessibility description, which is often the only
    # written account of a diagram. python-pptx still does the image FILES — markitdown
    # names images but does not write them, and the pack needs them on disk for triage.
    if _tool("markitdown"):
        md = _run(["markitdown", src])
        if md:
            body = re.sub(r"<!--\s*Slide number:\s*(\d+)\s*-->", r"## Slide \1:", md)
            by_slide = {}
            name_map = {}
            for i, slide in enumerate(prs.slides, 1):
                for shape in _pptx_shapes(slide.shapes):
                    if getattr(shape, "shape_type", None) == 13:
                        try:
                            name = f"s{i}_{shape.shape_id}.{shape.image.ext}"
                            if coll.add(shape.image.blob, i, name):
                                by_slide.setdefault(i, []).append(name)
                                key = re.sub(r"\s+", "", getattr(shape, "name", "") or "")
                                if key:
                                    name_map.setdefault(key, name)
                        except Exception:
                            continue
            # markitdown emits `![alt](Picture65.jpg)` for pictures it names but never writes
            # to disk — links that resolve nowhere (21 dead on the real GNI deck). Rewrite
            # each to a placement token for the file we DID extract, keeping markitdown's alt
            # text: it is often the only written account of the figure. Unmatchable links are
            # reduced to their alt text — a dead link helps no one.
            placed = set()

            def _fix_md_img(m):
                alt, href = m.group(1), m.group(2)
                key = re.sub(r"\s+", "", os.path.splitext(os.path.basename(href))[0])
                our = name_map.get(key)
                if not our:
                    return alt
                placed.add(our)
                safe_alt = alt.replace("|", "/").replace("--", "—")
                return IMG_TOKEN % (our + "|" + safe_alt) if safe_alt.strip() else IMG_TOKEN % our

            body = re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", _fix_md_img, body)
            needed = _pptx_diagram_slides(prs)
            snaps = _pptx_slide_snapshots(src, needed, coll)
            for i, name in snaps.items():
                by_slide.setdefault(i, []).append(name)
            # markitdown names pictures but writes no files, so it emits no link either.
            # Park each slide's remaining images (backgrounds markitdown skipped, slide
            # snapshots) at the end of that slide's block, not at the end of the deck: a
            # diagram three slides away from its narrative is a diagram nobody can caption.
            if by_slide:
                parts = re.split(r"(?m)^(## Slide (\d+):)$", body)
                rebuilt = [parts[0]]
                for k in range(1, len(parts), 3):
                    chunk = parts[k + 2]
                    toks = "".join("\n" + IMG_TOKEN % n for n in by_slide.get(int(parts[k + 1]), [])
                                   if n not in placed)
                    rebuilt.append(parts[k] + chunk.rstrip("\n") + toks + "\n")
                body = "".join(rebuilt)
            head = header(src, f"> **Slides:** {len(prs.slides)}  \n> **Extractor:** markitdown  ")
            if len(snaps) < len(needed):
                body += _snapshot_warning(len(needed) - len(snaps))
            tail = img_section(coll, len(prs.slides))
            return head + place_inline_images(body, coll) + tail, None

    out = [header(src, f"> **Slides:** {len(prs.slides)}  \n> **Extractor:** python-pptx (fallback — install markitdown for richer shape coverage + image alt-text)  ")]
    needed = _pptx_diagram_slides(prs)
    snaps = _pptx_slide_snapshots(src, needed, coll)
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}:\n")
        for shape in _pptx_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip() + "\n")
            if shape.has_table:
                tbl = shape.table
                rows = [[c.text.strip() for c in r.cells] for r in tbl.rows]
                if rows:
                    out.append("| " + " | ".join(rows[0]) + " |")
                    out.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                    for r in rows[1:]:
                        out.append("| " + " | ".join(r) + " |")
                    out.append("")
            if shape.shape_type == 13:  # picture
                try:
                    image = shape.image
                    name = f"s{i}_{shape.shape_id}.{image.ext}"
                    if coll.add(image.blob, i, name):
                        out.append(IMG_TOKEN % name)
                except Exception:
                    continue
        if i in snaps:
            out.append(IMG_TOKEN % snaps[i])
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    out.append(f"**Speaker notes:**\n\n{notes}\n")
        except Exception:
            pass
    if len(snaps) < len(needed):
        out.append(_snapshot_warning(len(needed) - len(snaps)))
    tail = img_section(coll, len(prs.slides))
    return place_inline_images("\n".join(out), coll) + tail, None


def convert_docx(src, images_dir, link_base=None):
    # Pandoc owns the extraction. It keeps document order (tables stay in their clause), real
    # heading levels, lists, footnotes and nested tables — all things a hand-rolled python-docx
    # walk gets wrong, and all things pandoc settled years ago. We add only what pandoc has no
    # reason to know about: the provenance header and the citable section numbering.
    if _tool("pandoc"):
        # --extract-media makes pandoc write the embedded images out; without it a docx's
        # diagrams are silently dropped, which no amount of good text extraction makes up for.
        # Pandoc has no view on which of them matter, so its output goes through the same
        # decorative filter as every other format.
        import shutil
        import tempfile
        tmp = tempfile.mkdtemp(prefix="engos-media-")
        md = _run(["pandoc", "-t", "gfm", "--wrap=none", f"--extract-media={tmp}", src])
        if md:
            coll = ImageCollector(images_dir, link_base, img_prefix(src))
            for root, _dirs, files in os.walk(tmp):
                for fn in sorted(files):
                    src_img = os.path.join(root, fn)
                    try:
                        with open(src_img, "rb") as fh:
                            coll.add(fh.read(), fn, fn)
                    except OSError:
                        continue
            shutil.rmtree(tmp, ignore_errors=True)
            tail = img_section(coll)                 # finalise() before any link_for() call

            # Pandoc points every image at the temp dir we just deleted. Repoint the survivors
            # at the pack's images dir and drop the links to images the filter removed, so the
            # markdown carries no dead references.
            #
            # BOTH syntaxes, not just markdown. Pandoc emits raw HTML `<img src=... style=.../>`
            # whenever the Word image carries an explicit size — which is nearly always. The
            # first version only rewrote `![](…)`, so on the real GNI tender pack every single
            # figure in every converted docx pointed into a temp directory that no longer
            # existed. The failure was invisible: the text conversion "succeeded".
            def _fix_md(m):
                rel = coll.link_for(m.group(2))
                return img_block(rel, m.group(1)) if rel else ""

            def _fix_html(m):
                tag = m.group(0)
                srcm = re.search(r'src\s*=\s*["\']([^"\']+)["\']', tag)
                if not srcm:
                    return ""
                rel = coll.link_for(srcm.group(1))
                if not rel:
                    return ""
                altm = re.search(r'alt\s*=\s*["\']([^"\']*)["\']', tag)
                return img_block(rel, altm.group(1) if altm else "")

            md = re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", _fix_md, md)
            md = re.sub(r"<img\b[^>]*/?>", _fix_html, md, flags=re.I)
            md = flatten_internal_links(md)
            body, n = number_headings(md, "Section")
            head = header(src, f"> **Sections:** {n}  \n> **Extractor:** pandoc  ")
            return head + body + tail, None

    # Fallback: no pandoc on this machine. Order is lost (python-docx exposes paragraphs and
    # tables as two separate sequences) — the header says so rather than pretending otherwise.
    try:
        import docx  # python-docx
    except ImportError:
        return None, "no pandoc, and python-docx not installed — install either: brew install pandoc | pip install python-docx (PEP-668: add --user --break-system-packages)"
    d = docx.Document(src)
    # A .docx has no stable page concept — pagination is decided by the renderer, so a "page"
    # anchor would not survive a font change. Anchor on SECTIONS instead (headings, or one
    # implicit section for a heading-less doc) so a claim is still citable as `file.md §Section N`.
    body, section, n = [], [], 0
    for para in d.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if "heading" in (para.style.name or "").lower():
            if section:
                body.append("\n".join(section))
            n += 1
            section = [f"## Section {n}: {text}\n"]
        else:
            if not section:                       # text before the first heading
                n += 1
                section = [f"## Section {n}:\n"]
            section.append(text + "\n")
    if section:
        body.append("\n".join(section))

    out = [header(src, f"> **Sections:** {n}  \n> **Extractor:** python-docx (fallback — tables are appended after the text, not in document order; install pandoc for faithful order)  ")]
    out.extend(body)
    for t, table in enumerate(d.tables, 1):
        rows = [[c.text.strip() for c in r.cells] for r in table.rows]
        if rows:
            out.append(f"\n**Table {t}:**\n")
            out.append("| " + " | ".join(rows[0]) + " |")
            out.append("| " + " | ".join("---" for _ in rows[0]) + " |")
            for r in rows[1:]:
                out.append("| " + " | ".join(r) + " |")
            out.append("")
    return "\n".join(out), None


def convert_xlsx(src, images_dir, link_base=None):
    try:
        import openpyxl
    except ImportError:
        return None, "openpyxl not installed — run: pip install openpyxl (or use the xlsx skill; PEP-668: add --user --break-system-packages)"
    wb = openpyxl.load_workbook(src, data_only=True, read_only=True)
    out = [header(src, f"> **Sheets:** {len(wb.sheetnames)}  ")]
    for name in wb.sheetnames:
        ws = wb[name]
        out.append(f"## Sheet: {name}\n")
        rows = list(ws.iter_rows(values_only=True))
        rows = [r for r in rows if any(c is not None for c in r)]
        if not rows:
            out.append("_(empty)_\n")
            continue
        width = max(len(r) for r in rows)
        def cell(x):
            return "" if x is None else str(x).replace("|", "\\|").replace("\n", " ")
        out.append("| " + " | ".join(cell(c) for c in (list(rows[0]) + [""] * (width - len(rows[0])))) + " |")
        out.append("| " + " | ".join("---" for _ in range(width)) + " |")
        for r in rows[1:]:
            padded = list(r) + [""] * (width - len(r))
            out.append("| " + " | ".join(cell(c) for c in padded) + " |")
        out.append("")
    return "\n".join(out), None


def convert_csv(src, images_dir, link_base=None):
    import csv
    out = [header(src)]
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        rows = list(csv.reader(f))
    rows = [r for r in rows if any(c.strip() for c in r)]
    if rows:
        out.append("| " + " | ".join(rows[0]) + " |")
        out.append("| " + " | ".join("---" for _ in rows[0]) + " |")
        for r in rows[1:]:
            out.append("| " + " | ".join(r) + " |")
    return "\n".join(out), None


def convert_html(src, images_dir, link_base=None):
    """A saved web page — the research lane's most common source, and the one with no file format.

    `defuddle` owns the de-cluttering, the way pandoc owns document order. A page saved from a
    corporate site is 90% navigation, cookie wall, share widgets and script tags; pandoc converts
    all of it faithfully into noise, which is worse than losing it. Defuddle is the same tool the
    `defuddle` skill reads pages with, so this is reuse, not a second implementation. We add only
    the provenance header and the citable section numbering.

    Images are remote here, not embedded, so there is nothing to extract to disk. Absolute links
    survive because they still resolve; relative ones are dropped, because a page saved without
    its origin turns `/content/dam/x.png` into a reference that points nowhere — the same dead-link
    failure convert_docx fixes for embedded media.
    """
    if _tool("defuddle"):
        md = _run(["defuddle", "parse", src, "-m"])
        extractor = "defuddle (page furniture removed)"
    else:
        md, extractor = None, None
    if not md and _tool("pandoc"):
        md = _run(["pandoc", "-f", "html", "-t", "gfm", "--wrap=none", src])
        extractor = ("pandoc (fallback — navigation, cookie banners and share widgets are "
                     "converted along with the article; install defuddle for a clean read)")
    if not md:
        return None, ("no defuddle and no pandoc — install either: "
                      "npm i -g defuddle-cli | brew install pandoc")

    kept, dropped = [0], [0]

    def _keep_absolute(m):
        url = m.group(2).strip()
        if url.lower().startswith(("http://", "https://", "data:")):
            kept[0] += 1
            return m.group(0)
        dropped[0] += 1
        return ""

    md = re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", _keep_absolute, md)

    # Unglue headings fused onto the previous paragraph. Web markup routinely closes a block
    # without a break, so the converter emits `...scales as your needs evolve.#### Zora AI for
    # Finance`. A heading that is not at line start is not a heading to any markdown parser, so
    # it silently skips the anchor numbering below — and a section with no anchor cannot be
    # cited, which is the whole point of ingesting it. Require a non-space before the hashes so
    # a legitimate mid-sentence `# ` is left alone.
    md = re.sub(r"([^\s])(#{1,6} )", r"\1\n\n\2", md)

    md = flatten_internal_links(md)
    body, n = number_headings(md, "Section")
    note = f"> **Images:** {kept[0]} remote link(s) kept, {dropped[0]} relative link(s) dropped  \n"
    return header(src, f"> **Sections:** {n}  \n{note}> **Extractor:** {extractor}  ") + body, None


def convert_image(src, images_dir, link_base=None):
    rel = os.path.relpath(src, link_base or os.path.dirname(images_dir)) if images_dir else src
    body = header(src) + (
        f"## Image\n\n`[{fc.UNTRIAGED}]` ![{os.path.basename(src)}]({rel})\n\n"
        "_This is an image source — OCR it inline (agent + vision), then retag `[ocr-done]`._\n"
    )
    return body, None


DISPATCH = {
    ".pdf": convert_pdf, ".pptx": convert_pptx, ".ppsx": convert_pptx,
    ".docx": convert_docx, ".xlsx": convert_xlsx, ".csv": convert_csv,
    ".png": convert_image, ".jpg": convert_image, ".jpeg": convert_image,
    ".html": convert_html, ".htm": convert_html,
}


def scan(root):
    """What showed up in the tree that hasn't been dealt with yet.

    Asking the user to type a path for every arriving document is asking them to do what a
    directory listing can decide. Two kinds of "not dealt with", because the two kinds of
    material have different destinations:

      INGEST  — sourced material (given to us or found by us, plus the tender pack). Un-ingested
        when nothing named after it exists under the matching `_md/` — where "named after it"
        is a same-stem md OR a manifest row mapping the source to a slugged md (the documented
        `_md/<NN_topic>/<slug>.md` convention).
      CONVERT — our OWN reusable assets under `01_pursuit/_shared/<kind>/` with no markdown
        under `_shared/_md/`. Assets get the SAME md-first treatment as sources: analysing
        binaries ad-hoc is what the pack forbids everywhere else ("analysing from the raw PDF
        is how citations get invented" — it is also how a €250k reference value gets misread
        on submission week).
      INDEX   — assets with no row in `firm_assets.md`. An asset with no row is invisible to
        the bid.
      RE-INDEX — assets whose markdown changed AFTER their row was written. A row drafted from
        a conversion whose figures were still untriaged says what the prose said and nothing
        the diagrams proved; on the real AIB case study that was the difference between
        "discovery-phase method" and a Teradata→VantageCloud interim-state transition, a
        2024-2027 quarterly roadmap and a phased code-conversion approach. Nothing noticed,
        because a stale row looks exactly like a considered one. Same mtime comparison
        `rule_estimate_snapshot_fresh` uses for a generated snapshot against its source.

    Returns (to_ingest, to_index, to_convert, to_reindex), each [(area_label, path)].
    """
    to_ingest = []
    areas = [(p, p / "_md") for p in root.glob("_sources/*") if p.is_dir()]
    areas += [(p, p / "_md") for p in root.glob("01_pursuit/*/1_received") if p.is_dir()]
    for area, md_dir in areas:
        converted = {p.stem.lower() for p in md_dir.rglob("*.md") if p.name != "README.md"} \
            if md_dir.exists() else set()
        converted |= manifest_aliases(md_dir)
        for src in sorted(area.rglob("*")):
            if not src.is_file() or src.suffix.lower() not in DISPATCH:
                continue
            if md_dir in src.parents:            # already inside _md/ — an extracted image
                continue
            if src.stem.lower() not in converted:
                to_ingest.append((str(area.relative_to(root)), src))

    to_index, to_convert, to_reindex = [], [], []
    for shared in root.glob("01_pursuit/_shared"):
        index_file = shared / "firm_assets.md"
        index_text = (index_file.read_text(encoding="utf-8", errors="replace")
                      if index_file.exists() else "")
        index_mtime = index_file.stat().st_mtime if index_file.exists() else None
        # A path mentioned under "Gaps — what we do NOT hold" is not handled evidence.
        # Keep the same boundary the lint uses when deciding which A-nnn ids are real assets.
        index_text = re.split(r"^#{1,3}\s+.*gap.*$", index_text,
                              flags=re.M | re.I)[0]
        # A basename substring is not an identity. On the real GNI asset tree,
        # `case_studies/others/Quals.pdf` disappeared because the index already contained
        # `case_studies/1_General_DataMod Quals.pdf`. Read exact, backticked paths instead.
        # The same mechanism lets the index mark editable companions and build helpers as
        # handled without pretending each one is a separate A-nnn evidence asset.
        indexed = set()
        primaries = set()
        for line in index_text.splitlines():
            tokens = re.findall(r"`([^`\n]+)`", line)
            # a row keyed by an A-nnn id is an Index row; its FIRST backticked path is the
            # conversion target — anything after it is a companion (one logical asset = one
            # conversion)
            is_row = bool(re.match(r"\s*\|?\s*A-\d{3}\b", line))
            for j, token in enumerate(tokens):
                rel = token.strip().replace("\\", "/")
                prefix = "01_pursuit/_shared/"
                if rel.startswith(prefix):
                    rel = rel[len(prefix):]
                if rel and not rel.startswith(("/", "../")) and "<" not in rel:
                    rel = pathlib.PurePosixPath(rel).as_posix()
                    indexed.add(rel)
                    if is_row and j == 0:
                        primaries.add(rel)
        # Walk ALL of _shared/, not a whitelist of kinds. The scaffold plants six, and the
        # folder's own README tells the user to add a new one rather than force a bad fit —
        # so a whitelist makes exactly the assets someone thought about hardest invisible.
        md_dir = shared / "_md"
        md_by_stem = {p.stem.lower(): p for p in md_dir.rglob("*.md")
                      if p.name not in ("README.md",)} if md_dir.exists() else {}
        converted = set(md_by_stem) | manifest_aliases(md_dir)
        for asset in sorted(shared.rglob("*")):
            if not asset.is_file() or asset.name.startswith("."):
                continue
            if asset.parent.name == "_md" or (shared / "_md") in asset.parents:
                continue                          # conversions, not assets
            if asset.name in ("README.md", "firm_assets.md"):
                continue
            # any file type — a rate card or a CV is an asset whether or not we can convert it
            rel = asset.relative_to(shared).as_posix()
            if rel not in indexed:
                where = asset.parent.relative_to(root)
                to_index.append((str(where), asset))
            elif index_mtime is not None and rel in primaries:
                # indexed, but is the row still true? Triage rewrites the md long after the
                # row was drafted, and everything a figure proves arrives in that rewrite.
                # Only the row's PRIMARY file is asked — a companion that happens to share a
                # stem (`X.drawio` beside `X.pptx`) is the same logical asset, not a second
                # row to redraft.
                md = md_by_stem.get(asset.stem.lower())
                if md and md.stat().st_mtime > index_mtime:
                    to_reindex.append((str(asset.parent.relative_to(root)), asset))
            if asset.suffix.lower() in DISPATCH and asset.stem.lower() not in converted:
                # one logical asset = one conversion: the row's FIRST path (its primary) or a
                # file nothing mentions (a new arrival). Companions and §1b helpers named
                # elsewhere in the index need no md of their own.
                if rel in primaries or rel not in indexed:
                    where = asset.parent.relative_to(root)
                    to_convert.append((str(where), asset))
    return to_ingest, to_index, to_convert, to_reindex


def manifest_aliases(md_dir):
    """Source stems the pack's manifest maps to an md of a different name.

    Same-stem matching covers `_md/<source>.md`; the manifest row covers the skill's own
    documented topic/slug convention (`_md/01_rft/rft_main.md` ← `26-002 - ....docx`).
    Without this the scan re-reports a converted document as waiting forever — the report
    everyone trusts crying wolf. Found on the 2026-07-29 mini GNI e2e.
    """
    aliases = set()
    readme = md_dir / "README.md"
    if readme.exists():
        for line in readme.read_text(encoding="utf-8", errors="replace").splitlines():
            m = re.match(r"\s*-\s+`([^`]+)`\s+—\s+converted from\s+`([^`]+)`", line)
            if m:
                aliases.add(pathlib.PurePosixPath(m.group(2)).stem.lower())
    return aliases


def update_manifest(out_path, src):
    """Upsert this conversion's provenance row in the pack's README.md manifest.

    Only fires when the output lands under a `_md/` pack directory — anywhere else there is no
    manifest convention to maintain. eng_lint's manifest-complete rule errors on a converted
    MD with no row, so if the converter doesn't write the row, every conversion lands the
    user in a lint error they can only clear by hand. Found on the real GNI tender pack:
    four documents converted, no manifest anywhere, and the error had no automated path out.

    The pack's `_md/` may be an ANCESTOR rather than the parent: the skill's documented
    convention is `_md/<NN_topic>/<slug>.md`, and requiring the parent to be `_md` itself
    meant the row was silently never written on the documented path — while stdout claimed
    it was (mini GNI e2e, 2026-07-29). Returns True when a row was written.

    Row format matches what rule_manifest_complete greps for: the .md name in backticks.
    """
    pack = pathlib.Path(out_path).parent
    while pack.name != "_md" and pack.parent != pack:
        pack = pack.parent
    if pack.name != "_md":
        return False
    readme = pack / "README.md"
    row = (f"- `{os.path.basename(out_path)}` — converted from "
           f"`{os.path.basename(src)}` ({_dt.date.today().isoformat()})")
    if not readme.exists():
        readme.write_text(
            "# Conversion manifest\n\n"
            "One row per converted document: which source file it came from and when. "
            "Written by convert_source.py on each conversion; keep rows in sync on renames.\n\n"
            + row + "\n", encoding="utf-8")
        return True
    text = readme.read_text(encoding="utf-8", errors="replace")
    name = os.path.basename(out_path)
    if f"`{name}`" in text:                      # re-conversion: refresh the row, don't duplicate
        text = re.sub(rf"^- `{re.escape(name)}`.*$", row, text, flags=re.M)
    else:
        text = text.rstrip("\n") + "\n" + row + "\n"
    readme.write_text(text, encoding="utf-8")
    return True


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("source", nargs="?",
                    help="the document to convert. Omit it with --scan to find what is waiting.")
    ap.add_argument("--scan", metavar="ROOT", nargs="?", const=".",
                    help="list source files in the engagement tree with no markdown yet, "
                         "and exit. Answers 'what arrived that I haven't ingested?'")
    ap.add_argument("--out", help="output .md path (default: source with .md next to it)")
    ap.add_argument("--images-dir", help="dir for extracted images (default: <out_dir>/images)")
    ap.add_argument("--source-url", help="the URL this document was downloaded from — recorded "
                                         "in the provenance header so a claim can cite something "
                                         "retrievable, not just a local path")
    args = ap.parse_args()

    global SOURCE_URL
    SOURCE_URL = args.source_url

    if args.scan is not None:
        root = pathlib.Path(args.scan).resolve()
        to_ingest, to_index, to_convert, to_reindex = scan(root)
        if not to_ingest and not to_index and not to_convert and not to_reindex:
            print("  nothing waiting — every source file has markdown, every asset has a row "
                  "and a conversion, and every row is newer than the markdown it reads")
            return 0

        def show(items, heading, follow_up):
            if not items:
                return
            print(f"  {len(items)} {heading}:\n")
            area = None
            for label, path in items:
                if label != area:
                    area, _ = label, print(f"  {label}/")
                print(f"      {path.relative_to(root / label)}")
            print(f"\n  → {follow_up}\n")

        show(to_ingest, "document(s) waiting to be INGESTED",
             "convert each, then update that bucket's canonical pair (eng-update-canonical).")
        show(to_convert, "asset(s) of OURS waiting to be CONVERTED to `_shared/_md/`",
             "eng-index-assets step 0 converts before indexing — the md layer is what analysis "
             "and page-level citation read; the binary is the evidence, not the working text.")
        show(to_index, "asset(s) of OURS waiting to be INDEXED",
             "eng-index-assets — what each proves, its date, whether it is in-window. "
             "An asset with no row is invisible to the bid.")
        show(to_reindex, "asset(s) waiting to be RE-INDEXED — markdown changed after the row",
             "eng-index-assets step 2, these rows only — image triage puts what the diagrams "
             "prove into the markdown, and a row drafted before it says what the prose said "
             "and nothing the figures showed.")
        return 0

    if not args.source:
        ap.error("give a document to convert, or --scan to list what is waiting")

    src = os.path.abspath(args.source)
    if not os.path.exists(src):
        print(f"ERROR: source not found: {src}", file=sys.stderr)
        return 2
    ext = os.path.splitext(src)[1].lower()
    fn = DISPATCH.get(ext)
    if not fn:
        print(f"ERROR: unsupported type '{ext}'. Supported: {', '.join(sorted(DISPATCH))}", file=sys.stderr)
        return 2

    out_path = os.path.abspath(args.out) if args.out else os.path.splitext(src)[0] + ".md"
    images_dir = os.path.abspath(args.images_dir) if args.images_dir else os.path.join(os.path.dirname(out_path), "images")
    os.makedirs(images_dir, exist_ok=True)

    body, err = fn(src, images_dir, os.path.dirname(out_path))
    if err:
        print(f"ERROR: {err}", file=sys.stderr)
        return 3
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(body)
    wrote_row = update_manifest(out_path, src)
    print(f"Wrote {out_path}")
    print(f"md5(source) = {md5(src)}")
    note = ("(the manifest row was written)" if wrote_row
            else "(output is outside a `_md/` pack — no manifest row to write)")
    print("Next: triage extracted images, OCR any [uncertain] ones inline "
          f"{note}, then run `eng-update-canonical`.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
